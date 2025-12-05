"""
Workflow service that loads embeddings, searches topics, and generates PDF summaries.
"""
import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime
from openai import OpenAI
from dotenv import load_dotenv

from powerpoint_maker_ddb.service.pdf_service import PDFService

# Load environment variables
load_dotenv()


class WorkflowService:
    """Service to orchestrate embedding search and PowerPoint generation workflow."""
    
    def __init__(self, vector_file: str = "src/powerpoint_maker_ddb/service/vectors.pkl", output_folder: str = "src/powerpoint"):
        """
        Initialize the workflow service.
        
        Args:
            vector_file: Path to vector storage file
            output_folder: Folder to save output PowerPoint files
        """
        self.pdf_service = PDFService(vector_file=vector_file)
        self.client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        self.output_folder = Path(output_folder)
        self.output_folder.mkdir(parents=True, exist_ok=True)
    
    def load_embeddings(self) -> Dict[str, Any]:
        """
        Load embeddings from storage.
        
        Returns:
            Dictionary with chunks, embeddings, and metadata
        """
        print("Loading embeddings...")
        vector_data = self.pdf_service.load_vectors()
        
        if len(vector_data["chunks"]) == 0:
            raise ValueError("No embeddings found. Please process PDFs first using PDFService.process_pdfs()")
        
        print(f"✓ Loaded {len(vector_data['chunks'])} chunks with embeddings")
        return vector_data
    
    def discover_topics(self, vector_data: Optional[Dict[str, Any]] = None, num_topics: int = 5, sample_size: int = 20) -> List[str]:
        """
        Automatically discover topics from the PDF content using AI analysis.
        
        Args:
            vector_data: Optional pre-loaded vector data
            num_topics: Number of topics to discover
            sample_size: Number of chunks to sample for topic discovery
            
        Returns:
            List of discovered topic strings
        """
        if vector_data is None:
            vector_data = self.load_embeddings()
        
        chunks = vector_data["chunks"]
        
        if len(chunks) == 0:
            raise ValueError("No chunks available for topic discovery")
        
        print(f"\nDiscovering topics from PDF content...")
        
        # Sample chunks for analysis
        if len(chunks) <= sample_size:
            sample_chunks = chunks
        else:
            step = len(chunks) // sample_size
            sample_chunks = [chunks[i] for i in range(0, len(chunks), step)][:sample_size]
        
        # Combine sample chunks for analysis
        sample_text = "\n\n---\n\n".join([f"[Section {i+1}]\n{chunk}" for i, chunk in enumerate(sample_chunks)])
        
        # Limit the text length to avoid token limits
        max_text_length = 8000
        if len(sample_text) > max_text_length:
            sample_text = sample_text[:max_text_length] + "..."
        
        # Create prompt for topic discovery
        prompt = f"""Analyze the following text content extracted from PDF documents and identify the {num_topics} main topics or themes discussed.

Text content:
{sample_text}

Please identify {num_topics} distinct main topics or themes. For each topic, provide a clear, concise topic name (2-5 words) that accurately represents the content.

Format your response as a simple list, one topic per line, like this:
Topic 1
Topic 2
Topic 3
...

Do not include numbers, bullets, or any other formatting - just the topic names, one per line."""

        try:
            response = self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant that analyzes documents and identifies main topics and themes accurately."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=300,
                temperature=0.5
            )
            
            # Parse the response to extract topics
            response_text = response.choices[0].message.content.strip()
            
            # Split by lines and clean up
            topics = []
            for line in response_text.split('\n'):
                line = line.strip()
                if line:
                    # Remove leading numbers, dots, dashes, etc.
                    cleaned = line.lstrip('0123456789.-) ').strip()
                    if cleaned and len(cleaned) > 2:
                        topics.append(cleaned)
            
            # Limit to requested number of topics
            topics = topics[:num_topics]
            
            if not topics:
                raise ValueError("No topics could be discovered from the content")
            
            print(f"  ✓ Discovered {len(topics)} topic(s):")
            for i, topic in enumerate(topics, 1):
                print(f"    {i}. {topic}")
            
            return topics
            
        except Exception as e:
            raise Exception(f"Error discovering topics: {str(e)}")
    
    def search_topics(self, topics: List[str], top_k: int = 5, vector_data: Optional[Dict[str, Any]] = None) -> Dict[str, List[Dict[str, Any]]]:
        """
        Search for each topic in the vector database.
        
        Args:
            topics: List of topic strings to search for
            top_k: Number of top results per topic
            vector_data: Optional pre-loaded vector data
            
        Returns:
            Dictionary mapping each topic to its search results
        """
        if vector_data is None:
            vector_data = self.load_embeddings()
        
        print(f"\nSearching for {len(topics)} topic(s)...")
        topic_results = {}
        
        for topic in topics:
            print(f"  Searching: '{topic}'")
            results = self.pdf_service.search_vectors(topic, top_k=top_k, vector_data=vector_data)
            topic_results[topic] = results
            print(f"    Found {len(results)} relevant chunks")
        
        return topic_results
    
    def summarize_information(self, topic: str, chunks: List[str], max_length: int = 500) -> str:
        """
        Summarize information from relevant chunks for a topic.
        
        Args:
            topic: The topic being summarized
            chunks: List of relevant text chunks
            max_length: Maximum length of summary in characters
            
        Returns:
            Summarized text
        """
        if not chunks:
            return f"No relevant information found for topic: {topic}"
        
        # Combine chunks with context
        combined_text = "\n\n".join([f"[Chunk {i+1}]\n{chunk}" for i, chunk in enumerate(chunks)])
        
        # Create prompt for summarization
        prompt = f"""Please provide a concise summary of the following information related to the topic: "{topic}"

Information:
{combined_text}

Please summarize the key points about "{topic}" in a clear and structured way. Keep the summary under {max_length} characters."""

        try:
            response = self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant that summarizes information clearly and concisely."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=500,
                temperature=0.3
            )
            
            summary = response.choices[0].message.content.strip()
            return summary
        except Exception as e:
            raise Exception(f"Error summarizing information: {str(e)}")
    
    def generate_powerpoint(self, topic_summaries: Dict[str, str], output_filename: Optional[str] = None) -> Path:
        """
        Generate a PowerPoint presentation with summarized information for each topic.
        
        Args:
            topic_summaries: Dictionary mapping topics to their summaries
            output_filename: Optional custom output filename
            
        Returns:
            Path to the generated PowerPoint file
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError("python-pptx is required. Install it with: pip install python-pptx")
        
        # Generate output filename with timestamp
        if output_filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = f"topic_summaries_{timestamp}.pptx"
        
        if not output_filename.endswith('.pptx'):
            output_filename += '.pptx'
        
        output_path = self.output_folder / output_filename
        
        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Topic Summaries Report"
        subtitle.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Add a slide for each topic
        for i, (topic, summary) in enumerate(topic_summaries.items(), 1):
            # Use blank layout for more control
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add topic title
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = f"{i}. {topic}"
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.size = Pt(32)
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.LEFT
            
            # Add summary content
            content_top = Inches(1.8)
            content_height = Inches(5)
            content_box = slide.shapes.add_textbox(left, content_top, width, content_height)
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            # Split summary into paragraphs if it's long
            summary_lines = summary.split('\n')
            for j, line in enumerate(summary_lines):
                if j == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                p.text = line.strip()
                p.font.size = Pt(14)
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(6)
        
        # Save presentation
        prs.save(str(output_path))
        
        print(f"\n✓ PowerPoint generated: {output_path}")
        return output_path
    
    def run_workflow(self, topics: Optional[List[str]] = None, top_k: int = 5, output_filename: Optional[str] = None, num_topics: int = 5) -> Dict[str, Any]:
        """
        Run the complete workflow: load embeddings, discover/search topics, summarize, and generate PowerPoint.
        
        Args:
            topics: Optional list of topics to search and summarize. If None, topics will be auto-discovered.
            top_k: Number of top results to retrieve per topic
            output_filename: Optional custom output filename
            num_topics: Number of topics to discover if topics are not provided
            
        Returns:
            Dictionary with workflow results
        """
        print("=" * 60)
        print("Starting Workflow Service")
        print("=" * 60)
        
        # Step 1: Load embeddings
        vector_data = self.load_embeddings()
        
        # Step 2: Discover topics if not provided
        if topics is None:
            topics = self.discover_topics(vector_data=vector_data, num_topics=num_topics)
        else:
            print(f"\nUsing provided topics: {topics}")
        
        # Step 3: Search topics
        topic_results = self.search_topics(topics, top_k=top_k, vector_data=vector_data)
        
        # Step 4: Loop through topics and summarize
        print(f"\nSummarizing information for {len(topics)} topic(s)...")
        topic_summaries = {}
        
        for topic in topics:
            print(f"  Summarizing: '{topic}'")
            # Get chunks from search results
            chunks = [result["chunk"] for result in topic_results[topic]]
            
            # Summarize
            summary = self.summarize_information(topic, chunks)
            topic_summaries[topic] = summary
            print(f"    ✓ Summary created ({len(summary)} characters)")
        
        # Step 5: Generate PowerPoint
        print(f"\nGenerating PowerPoint with summaries...")
        pptx_path = self.generate_powerpoint(topic_summaries, output_filename)
        
        print("\n" + "=" * 60)
        print("Workflow completed successfully!")
        print("=" * 60)
        
        return {
            "status": "success",
            "topics_processed": len(topics),
            "topics": topics,
            "pptx_path": str(pptx_path),
            "summaries": topic_summaries
        }

